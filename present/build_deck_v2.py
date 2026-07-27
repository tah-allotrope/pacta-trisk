#!/usr/bin/env python3
"""Build the Bank-Prospect PACTA/TRISK Pitch Deck v2.

Follows the `present` skill procedure:
- set_text: clears frame, zeroes margins, single run, <a:buNone/> to suppress bullets
- delete_shape + add_picture at original coordinates for images
- Saves to a new output path (never overwrites the reference)
"""

import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

REPO_ROOT = Path(__file__).parent.parent
REF_PATH = REPO_ROOT / "present" / "ref" / "Allotrope GTB updates_Tara Vietnam_December 2025.pptx"
ASSETS_DIR = REPO_ROOT / "present" / "assets"
OUTPUT_PATH = REPO_ROOT / "present" / "bank_prospect_deck_v2.pptx"

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

SLIDES_TO_DROP = [13, 14, 17, 18, 19, 20, 21, 22, 23, 24]


def set_text(shape, text, size=None, bold=None, color=None, font=None, align=None):
    """Replace all text in a shape with a single run.

    Clears the text frame, zeroes margins, sets one run, and suppresses
    any inherited bullet by appending <a:buNone/> to the paragraph pPr.
    """
    tf = shape.text_frame
    tf.word_wrap = True

    for p in tf.paragraphs:
        p.clear()
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.clear()
        run = p.add_run()
        run.text = line

        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = RGBColor(*color)
        if font is not None:
            run.font.name = font
        if align is not None:
            p.alignment = align

        pPr = p._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        if pPr is None:
            pPr = etree.SubElement(
                p._p,
                "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr",
            )
            p._p.insert(0, pPr)

        buNone = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
        if buNone is None:
            etree.SubElement(
                pPr,
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone",
            )

    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def delete_shape(slide, shape):
    """Remove a shape from a slide."""
    sp = shape._element
    sp.getparent().remove(sp)


def add_image(slide, image_path, x, y, w, h):
    """Add an image at the given coordinates (inches)."""
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))


def clear_all_text(slide):
    """Set all text on a slide to empty, preserving images and decorative shapes."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            tf = sh.text_frame
            tf.clear()
            for p in tf.paragraphs:
                p.clear()
                run = p.add_run()
                run.text = ""


def clear_table(table_shape):
    """Clear all text from table cells."""
    table = table_shape.table
    for row in table.rows:
        for cell in row.cells:
            cell.text = ""
            tf = cell.text_frame
            tf.clear()
            for p in tf.paragraphs:
                p.clear()
                run = p.add_run()
                run.text = ""


def set_table_cells(table_shape, rows_data):
    """Populate table cells from a list of lists."""
    table = table_shape.table
    for r, row_data in enumerate(rows_data):
        if r >= len(table.rows):
            break
        for c, cell_text in enumerate(row_data):
            if c >= len(table.rows[r].cells):
                break
            cell = table.rows[r].cells[c]
            cell.text = cell_text


def drop_slides(prs, indices):
    """Remove slides by index in reverse order."""
    for idx in sorted(indices, reverse=True):
        if idx < len(prs.slides):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]


def build():
    print(f"Loading reference from {REF_PATH}")
    prs = Presentation(str(REF_PATH))
    print(f"  {len(prs.slides)} slides, {Emu(prs.slide_width).inches:.2f}x{Emu(prs.slide_height).inches:.2f} in")

    print(f"Dropping {len(SLIDES_TO_DROP)} slides: {SLIDES_TO_DROP}")
    drop_slides(prs, SLIDES_TO_DROP)
    print(f"  {len(prs.slides)} slides remaining")

    slides = list(prs.slides)
    # After dropping [13,14,17-24], remaining ref slides map to:
    # slides[0]=ref0(Cover) [1]=ref1(Agenda) [2]=ref2(Table) [3]=ref3(BIDV-25shapes)
    # [4]=ref4(Challenges) [5]=ref5(TCB-20shapes) [6]=ref6(LookingAhead) [7]=ref7(Contact)
    # [8]=ref8(SiteAnalysis1) [9]=ref9(SiteAnalysis1-detail) [10]=ref10(5-box)
    # [11]=ref11(BIDV-Scope) [12]=ref12(BIDV-Scope2) [13]=ref15(LookingAhead2) [14]=ref16(Questions)

    # --- Slide 0 (ref 0): Cover ---
    print("  Slide 0: Cover")
    s = slides[0]
    clear_all_text(s)
    set_text(s.shapes[2],
             "PACTA + TRISK: Transition-Risk Analytics\nfor Vietnamese Banks",
             size=28, bold=True, color=(0, 102, 102),
             font="Calibri", align=PP_ALIGN.CENTER)

    # --- Slide 1 (ref 1): Agenda ---
    print("  Slide 1: Agenda")
    s = slides[1]
    clear_all_text(s)
    set_text(s.shapes[1], "Agenda",
             size=28, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[2], "The Problem  |  Our Approach  |  Case Study  |  Next Steps",
             size=12, color=(102, 102, 102), font="Calibri")

    # --- Slide 2 (ref 2): Our Solution (table layout) ---
    print("  Slide 2: Our Solution")
    s = slides[2]
    clear_all_text(s)
    set_text(s.shapes[1],
             "Our Solution: Three-Layer Compliance & Risk Stack",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")
    clear_table(s.shapes[2])
    set_table_cells(s.shapes[2], [
        ["Layer", "Tool", "Purpose", "Output"],
        ["1. Measure", "PCAF", "GHG inventory", "Scope 1, 2, 3 emissions"],
        ["2. Align", "PACTA", "Portfolio alignment", "PDP8 / NDC / NZE gaps"],
        ["3. Stress", "TRISK", "Borrower stress", "NPV & PD impact"],
        ["4. Disclose", "IFRS S2", "Board reporting", "TCFD / ISSB outputs"],
        ["", "", "", ""],
    ])

    # --- Slide 3 (ref 3): PACTA (25 shapes) ---
    print("  Slide 3: PACTA")
    s = slides[3]
    clear_all_text(s)
    set_text(s.shapes[1],
             "PACTA: Portfolio Alignment Assessment",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[4],
             "Develop portfolio alignment reports using open-source R packages",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[5],
             "43 loans, 25 trillion VND synthetic portfolio across power, automotive, cement, steel, and coal mining sectors",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[6],
             "Market-share method for power and automotive",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[7],
             "SDA method for cement and steel emission intensity",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[8],
             "Custom PDP8 scenario from Vietnam Power Development Plan 8",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[9],
             "Company-level results enable borrower-specific engagement",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[10],
             "Scenario comparison: PDP8 vs IEA NZE vs Vietnam NDC",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[11],
             "Open-source R packages, no licensing fees",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[14],
             "1,500+ institutions globally",
             size=8, color=(102, 102, 102), font="Calibri")
    set_text(s.shapes[16],
             "Coal 28% above PDP8 target",
             size=9, bold=True, color=(180, 40, 40), font="Calibri")
    set_text(s.shapes[17],
             "Renewables 15% below 2030 trajectory",
             size=9, bold=True, color=(180, 40, 40), font="Calibri")
    set_text(s.shapes[24],
             "Time to first results: 1-3 months",
             size=8, color=(102, 102, 102), font="Calibri")

    img_path = ASSETS_DIR / "05_vn_power_techmix.png"
    if img_path.exists():
        add_image(s, img_path, 6.60, 1.17, 3.12, 2.40)

    # --- Slide 4 (ref 4): The Challenge (8 shapes) ---
    print("  Slide 4: The Challenge")
    s = slides[4]
    clear_all_text(s)
    set_text(s.shapes[1],
             "The Challenge: Decision 263 & Transition Risk",
             size=20, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[3],
             "Decision 263/QD-TTg mandates GHG quotas for thermal power, steel, and cement from 2025. Banks must assess borrower transition risk, but lack analytical tools to quantify portfolio alignment or borrower-level financial stress under Paris-aligned scenarios.",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[4],
             "Vietnamese banks need an integrated analytics stack that measures portfolio alignment with PDP8/NDC targets, quantifies borrower-level NPV and PD impact under transition scenarios, and produces disclosure-ready outputs for IFRS S2 and TCFD.",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[5], "THE CHALLENGE",
             size=9, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[6], "WHAT BANKS NEED",
             size=9, bold=True, color=(51, 51, 51), font="Calibri")

    # --- Slide 5 (ref 5): TRISK (20 shapes) ---
    print("  Slide 5: TRISK")
    s = slides[5]
    clear_all_text(s)
    set_text(s.shapes[1],
             "TRISK: Borrower-Level Transition Stress Testing",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[3],
             "Quantify borrower-level NPV and PD impact under transition scenarios",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[4],
             "Power: Coal-heavy borrowers show -12% to -18% NPV change",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[5],
             "Cement: BF/BOF mills face -8% to -14% NPV impact",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[6],
             "Steel: EAF routes outperform BF/BOF by 6-10pp",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[7],
             "Priority ranking: alignment gap + NPV change + exposure",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[8],
             "Sensitivity: shock year, discount rate, market passthrough",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[9],
             "Scenario Builder: drive five TRISK levers in real time",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[10],
             "Covers all three Decision 263 sectors",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[12],
             "Exportable: CSV, HTML reports, interactive dashboard",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[13],
             "Results update instantly, no recomputation needed",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[19],
             "Open-source R package (trisk.model) on CRAN",
             size=8, color=(102, 102, 102), font="Calibri")

    img_path = ASSETS_DIR / "trisk_power_npv_change.png"
    if img_path.exists():
        add_image(s, img_path, 6.56, 1.62, 3.12, 2.60)

    # --- Slide 6 (ref 6): Engagement Offer (4 shapes) ---
    print("  Slide 6: Engagement Offer")
    s = slides[6]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Engagement Offer & 2026 Roadmap",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[3],
             "Phase 1: Data intake and loanbook normalization (BYOL)\nPhase 2: PACTA portfolio alignment with real borrower data\nPhase 3: TRISK stress testing and borrower ranking\nPhase 4: IFRS S2 disclosure package and board-ready report\n\nTimeline: 6-9 months from data receipt to final deliverables.\nAll analytics are open-source with no licensing fees.",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[2],
             "We welcome the opportunity to discuss how this analytics stack can support your bank's transition-risk management and Decision 263 compliance.",
             size=10, color=(68, 68, 68), font="Calibri")

    # --- Slide 7 (ref 7): Contact (5 shapes) ---
    print("  Slide 7: Contact")
    s = slides[7]
    clear_all_text(s)
    set_text(s.shapes[4],
             "Tung Ho, tah@allotropepartners.com\nHang Tran, httt@allotropepartners.com\nAllotrope VC | GTB Vietnam",
             size=14, bold=False, color=(51, 51, 51), font="Calibri")

    # --- Slide 8 (ref 8): Case Study Power (3 shapes) ---
    print("  Slide 8: Case Study Power")
    s = slides[8]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Case Study: Vietnam Bank Demo - Power Sector Alignment",
             size=16, bold=True, color=(51, 51, 51), font="Calibri")

    table_shape = s.shapes[2]
    clear_table(table_shape)
    set_table_cells(table_shape, [
        ["Metric", "Value"],
        ["Portfolio", "43 loans, ~$950M USD synthetic"],
        ["Power Mix", "Coal 28%, Gas 12%, Hydro 10%, Solar 8%, Wind 5%"],
        ["PDP8 Target", "Coal phase-down, renewables 3x by 2030"],
        ["Gap", "Coal +28% above target; renewables -15% below"],
    ])

    img_path = ASSETS_DIR / "12_vn_alignment_overview.png"
    if img_path.exists():
        old_pic = s.shapes[1]
        delete_shape(s, old_pic)
        add_image(s, img_path, 5.65, 1.22, 4.04, 2.84)

    # --- Slide 9 (ref 9): Case Study TRISK (7 shapes) ---
    print("  Slide 9: Case Study TRISK")
    s = slides[9]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Case Study: TRISK Borrower Stress Rankings",
             size=16, bold=True, color=(51, 51, 51), font="Calibri")

    set_text(s.shapes[4], "Priority Score (Top 10)",
             size=9, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[5], "PD Change by Borrower",
             size=9, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[6],
             "TRISK ranks borrowers by financial stress severity. Coal-heavy names (Vinh Tan, Duyen Hai, Mong Duong) show highest NPV decline and PD increase. Renewable-focused borrowers (Trung Nam, BIM) benefit from transition tailwinds.\n\nIllustrative / synthetic data - not actual portfolio results.",
             size=9, color=(68, 68, 68), font="Calibri")

    for shape_idx, asset_name in sorted([(1, "trisk_power_priority_top10.png"),
                                   (2, "trisk_power_pd_change.png"),
                                   (3, "13_vn_coal_stranded_risk.png")], reverse=True):
        img_path = ASSETS_DIR / asset_name
        if img_path.exists():
            old_pic = s.shapes[shape_idx]
            x, y = Emu(old_pic.left).inches, Emu(old_pic.top).inches
            w, h = Emu(old_pic.width).inches, Emu(old_pic.height).inches
            delete_shape(s, old_pic)
            add_image(s, img_path, x, y, w, h)

    # --- Slide 10 (ref 10): Key Findings (27 shapes) ---
    print("  Slide 10: Key Findings")
    s = slides[10]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Case Study: Key Findings & Takeaways",
             size=16, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[12], "Coal Stranded Risk",
             size=10, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[13], "Renewables Gap",
             size=10, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[14], "Borrower Heterogeneity",
             size=10, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[15], "Decision 263 Readiness",
             size=10, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[16], "Interactive Exploration",
             size=10, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[17],
             "28% of portfolio coal capacity exceeds PDP8 phase-down trajectory. BOT plants face cliff-edge decline post-PPA expiry (~2035).",
             size=8, color=(255, 255, 255), font="Calibri")
    set_text(s.shapes[18],
             "Renewables buildout is 15% below 2030 target. Solar and wind capacity must triple to align with PDP8.",
             size=8, color=(255, 255, 255), font="Calibri")
    set_text(s.shapes[19],
             "NPV impact ranges from -18% (coal-heavy) to +8% (renewable-focused). One-size engagement fails.",
             size=8, color=(255, 255, 255), font="Calibri")
    set_text(s.shapes[20],
             "Three-layer stack (PCAF, PACTA, TRISK) addresses all Decision 263 requirements: inventory, quotas, and reduction plans.",
             size=8, color=(255, 255, 255), font="Calibri")
    set_text(s.shapes[21],
             "Live dashboard with Scenario Builder lets bank teams explore levers in real time. Downloadable HTML reports and CSV exports.",
             size=8, color=(255, 255, 255), font="Calibri")

    # --- Slide 11 (ref 11): Compliance Value (8 shapes) ---
    print("  Slide 11: Compliance Value")
    s = slides[11]
    clear_all_text(s)
    set_text(s.shapes[1],
             "Compliance Value: Decision 263 / TCFD / ISSB",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")

    set_text(s.shapes[5],
             "Cement & Steel SDA: Emission intensity vs. 2030 NDC targets",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[6],
             "Coal trajectory vs. PDP8 phase-down pathway",
             size=9, color=(68, 68, 68), font="Calibri")

    for shape_idx, asset_name in sorted([(2, "10_vn_cement_sda.png"),
                                   (3, "11_vn_steel_sda.png"),
                                   (4, "06_vn_coal_trajectory.png")], reverse=True):
        img_path = ASSETS_DIR / asset_name
        if img_path.exists():
            old_pic = s.shapes[shape_idx]
            x, y = Emu(old_pic.left).inches, Emu(old_pic.top).inches
            w, h = Emu(old_pic.width).inches, Emu(old_pic.height).inches
            delete_shape(s, old_pic)
            add_image(s, img_path, x, y, w, h)

    # --- Slide 12 (ref 12): Live Dashboard (10 shapes) ---
    print("  Slide 12: Live Dashboard")
    s = slides[12]
    clear_all_text(s)
    set_text(s.shapes[1],
             "Live Dashboard: Interactive Portfolio Exploration",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")

    set_text(s.shapes[6],
             "PACTA Alignment: Sector KPIs and interactive charts",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[7],
             "TRISK Risk: Borrower stress rankings with sensitivity controls",
             size=9, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[8],
             "Scenario Builder: Drive the five TRISK levers in real time",
             size=9, color=(68, 68, 68), font="Calibri")

    img_path = ASSETS_DIR / "03_vn_coverage_pie.png"
    if img_path.exists():
        old_pic = s.shapes[2]
        x, y = Emu(old_pic.left).inches, Emu(old_pic.top).inches
        w, h = Emu(old_pic.width).inches, Emu(old_pic.height).inches
        delete_shape(s, old_pic)
        add_image(s, img_path, x, y, w, h)

    for shape_idx in sorted([3, 4, 5], reverse=True):
        img_path = ASSETS_DIR / "12_vn_alignment_overview.png"
        if img_path.exists():
            old_pic = s.shapes[shape_idx]
            x, y = Emu(old_pic.left).inches, Emu(old_pic.top).inches
            w, h = Emu(old_pic.width).inches, Emu(old_pic.height).inches
            delete_shape(s, old_pic)
            add_image(s, img_path, x, y, w, h)

    # --- Slide 13 (ref 15): Next Steps (4 shapes) ---
    print("  Slide 13: Next Steps")
    s = slides[13]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Proposed Next Steps",
             size=18, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[3],
             "1. Share loanbook template for BYOL intake validation\n2. Identify Decision 263 borrower subset for priority analysis\n3. Schedule technical workshop with risk and ESG teams\n4. Define scenario preferences (PDP8, NDC, IEA NZE)\n5. Agree on reporting language and disclosure timeline",
             size=10, color=(68, 68, 68), font="Calibri")
    set_text(s.shapes[2],
             "We are ready to begin Phase 1 upon receipt of loanbook data.\nSynthetic demo available for immediate internal training.",
             size=10, color=(68, 68, 68), font="Calibri")

    # --- Slide 14 (ref 16): Questions (2 shapes) ---
    print("  Slide 14: Questions")
    s = slides[14]
    clear_all_text(s)
    set_text(s.shapes[0],
             "Questions & Discussion",
             size=24, bold=True, color=(51, 51, 51), font="Calibri")
    set_text(s.shapes[1],
             "We welcome questions on methodology, data requirements, implementation timeline, or any aspect of the PACTA + TRISK analytics stack.",
             size=12, color=(102, 102, 102), font="Calibri")

    print(f"Saving to {OUTPUT_PATH}")
    prs.save(str(OUTPUT_PATH))
    print("Done!")


if __name__ == "__main__":
    build()
