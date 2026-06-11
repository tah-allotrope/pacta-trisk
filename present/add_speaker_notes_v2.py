#!/usr/bin/env python3
"""Add speaker notes to the built Bank-Prospect PACTA/TRISK Pitch Deck."""

from pathlib import Path
from pptx import Presentation

REPO_ROOT = Path(__file__).parent.parent
DECK_PATH = REPO_ROOT / "present" / "bank_prospect_deck_v2.pptx"

SPEAKER_NOTES = {
    0: (
        "Welcome and introductions. This deck presents our PACTA + TRISK analytics stack "
        "for Vietnamese banks navigating Decision 263 and transition-risk compliance.\n\n"
        "Key message: We offer an integrated, open-source, no-licensing-fee solution that "
        "takes a bank from GHG inventory through portfolio alignment to borrower-level stress "
        "testing and IFRS S2 disclosure.\n\n"
        "Synthetic data disclaimer: All case-study results in this deck use illustrative synthetic "
        "data — not actual portfolio results."
    ),
    1: (
        "Walk through the four-part agenda:\n"
        "1. The Problem — Decision 263 mandates and the analytical gap\n"
        "2. Our Approach — Three-layer compliance stack (PCAF, PACTA, TRISK)\n"
        "3. Case Study — Vietnam bank demo with synthetic data\n"
        "4. Next Steps — How to engage and timeline"
    ),
    2: (
        "Frame the challenge: Decision 263/QĐ-TTg (signed July 2024) requires banks to:\n"
        "- Report GHG emissions from financed activities\n"
        "- Set reduction targets aligned with Vietnam's NDC (43.5% conditional by 2030)\n"
        "- Develop transition plans for high-carbon borrowers\n\n"
        "The problem: Most Vietnamese banks lack the analytical tools to quantify portfolio "
        "alignment or borrower-level financial stress under transition scenarios.\n\n"
        "Ask: 'How is your bank currently assessing transition risk for your loan portfolio?'"
    ),
    3: (
        "Present the four-layer stack:\n"
        "- Layer 1 PCAF: Measure — GHG inventory for Scope 1, 2, and 3 (financed and non-financed)\n"
        "- Layer 2 PACTA: Align — Portfolio alignment against PDP8, NDC, and IEA NZE scenarios\n"
        "- Layer 3 TRISK: Stress — Borrower-level NPV and PD impact under transition scenarios\n"
        "- Layer 4 IFRS S2: Disclose — TCFD/ISSB-compliant reporting outputs\n\n"
        "Key differentiator: All layers use open-source R packages (r2dii.*, pacta.loanbook, trisk.model) "
        "with no licensing fees. The stack is modular — banks can adopt one layer or all four."
    ),
    4: (
        "PACTA deep-dive:\n"
        "- Market-share method for power and automotive sectors\n"
        "- Sectoral Decarbonization Approach (SDA) for cement and steel\n"
        "- Fuzzy matching with manual review for Vietnamese company names\n"
        "- Custom PDP8 scenario interpolation from Vietnam's Power Development Plan 8\n\n"
        "Synthetic demo results: 43-loan portfolio (~$1B USD) shows coal capacity 28% above "
        "PDP8 target and renewables 15% below 2030 trajectory.\n\n"
        "Emphasize: Company-level results enable borrower-specific engagement, not just "
        "portfolio-level averages."
    ),
    5: (
        "TRISK deep-dive:\n"
        "- DCF-based NPV modeling for each borrower's physical assets\n"
        "- Merton-model PD estimation under transition stress\n"
        "- Five configurable levers: shock year, discount rate, risk-free rate, market passthrough, "
        "carbon-price family\n"
        "- Sensitivity analysis built into the Scenario Builder dashboard page\n\n"
        "Key finding: Coal-heavy borrowers show -12% to -18% NPV change under NZE-aligned stress, "
        "while renewable-focused borrowers benefit from transition tailwinds (+8% NPV).\n\n"
        "Emphasize: Results update instantly — no recomputation needed when levers change."
    ),
    6: (
        "Case study: Power sector alignment\n"
        "- Synthetic MCB portfolio: 43 loans across power, automotive, cement, steel, coal\n"
        "- Power mix: Coal 28%, Gas 12%, Hydro 10%, Solar 8%, Wind 5%\n"
        "- PDP8 target: Coal phase-down, renewables must triple by 2030\n"
        "- Gap: Coal capacity 28% above target; renewables 15% below\n\n"
        "Point to the alignment overview chart: Red bars = misaligned, green bars = aligned.\n\n"
        "Disclaimer reminder: This is illustrative synthetic data. Real bank results will differ "
        "based on actual loanbook composition."
    ),
    7: (
        "Case study: TRISK borrower stress rankings\n"
        "- Priority score combines alignment gap, NPV change, and exposure weight\n"
        "- Top stressed borrowers: Coal-heavy names (Vinh Tan, Duyen Hai, Mong Duong) show "
        "highest NPV decline and PD increase\n"
        "- Renewable-focused borrowers (Trung Nam, BIM) benefit from transition tailwinds\n\n"
        "Point to the three charts:\n"
        "1. Priority score ranking (top 10)\n"
        "2. PD change by borrower\n"
        "3. Coal stranded risk visualization\n\n"
        "Key insight: One-size engagement fails. Borrower-specific data enables targeted "
        "transition planning."
    ),
    8: (
        "Five key takeaways from the case study:\n"
        "1. Coal Stranded Risk: 28% of portfolio coal capacity exceeds PDP8 phase-down trajectory. "
        "BOT plants face cliff-edge decline post-PPA expiry (~2035).\n"
        "2. Renewables Gap: Solar and wind capacity must triple to align with PDP8.\n"
        "3. Borrower Heterogeneity: NPV impact ranges from -18% (coal) to +8% (renewables).\n"
        "4. Decision 263 Readiness: Three-layer stack addresses all requirements.\n"
        "5. Interactive Exploration: Live dashboard with Scenario Builder.\n\n"
        "Ask: 'Which of these findings resonates most with your current portfolio concerns?'"
    ),
    9: (
        "Compliance value proposition:\n"
        "- Decision 263: GHG inventory (PCAF), reduction quotas (PACTA alignment), transition "
        "plans (TRISK stress testing)\n"
        "- TCFD: Four-pillar disclosure (Governance, Strategy, Risk Management, Metrics & Targets)\n"
        "- ISSB IFRS S2: Climate-related financial disclosures\n\n"
        "Point to the three charts:\n"
        "- Cement SDA: Emission intensity vs. 2030 NDC targets\n"
        "- Steel SDA: BF/BOF vs. EAF route comparison\n"
        "- Coal trajectory: Phase-down pathway vs. current exposure\n\n"
        "Key message: Our stack produces board-ready outputs that satisfy all three frameworks "
        "simultaneously."
    ),
    10: (
        "Live dashboard demonstration:\n"
        "- 7-page Streamlit app deployed at pactavn.streamlit.app\n"
        "- PACTA Alignment page: Sector-level KPIs, interactive charts, downloadable CSVs\n"
        "- TRISK Risk page: Borrower stress rankings with sensitivity controls\n"
        "- Scenario Builder page: Drive the five TRISK levers in real time\n"
        "- Reports page: Inline HTML reports with download buttons\n"
        "- Methodology page: PACTA/TRISK framing with citations\n"
        "- Outputs page: Engagement letters and disclosure pack generation\n\n"
        "Offer to do a live demo if time permits. The dashboard is publicly accessible for "
        "exploration."
    ),
    11: (
        "Engagement offer:\n"
        "- Phase 1: BYOL (Bring Your Own Loanbook) — 2-4 weeks\n"
        "- Phase 2: PACTA alignment with real borrower data — 4-6 weeks\n"
        "- Phase 3: TRISK stress testing and borrower ranking — 3-4 weeks\n"
        "- Phase 4: IFRS S2 disclosure package and board-ready report — 2-3 weeks\n\n"
        "Total timeline: 6-9 months from data receipt to final deliverables.\n\n"
        "Emphasize: All analytics are open-source with no licensing fees. The only cost is "
        "our consulting time for data intake, validation, and report generation."
    ),
    12: (
        "Proposed next steps:\n"
        "1. Share loanbook template for BYOL intake validation (we provide the template)\n"
        "2. Identify Decision 263 borrower subset for priority analysis (top 20 by exposure)\n"
        "3. Schedule technical workshop with risk and ESG teams (half-day session)\n"
        "4. Define scenario preferences (PDP8, NDC, IEA NZE — we support all three)\n"
        "5. Agree on reporting language and disclosure timeline\n\n"
        "We are ready to begin Phase 1 upon receipt of loanbook data.\n\n"
        "Synthetic demo is available for immediate internal training — no data required."
    ),
    13: (
        "Open the floor for questions:\n"
        "- Methodology: PACTA matching, TRISK modeling, scenario selection\n"
        "- Data requirements: Loanbook format, ABCD sources, scenario files\n"
        "- Implementation timeline: Phase durations, dependencies, parallelization\n"
        "- Cost structure: Consulting fees vs. open-source licensing\n"
        "- Integration: How outputs feed into existing risk and ESG workflows\n\n"
        "Be prepared to discuss:\n"
        "- Vietnamese company name matching challenges (diacritics, VSIC codes)\n"
        "- BOT coal plant contractual lock-in (PPA expiry dates)\n"
        "- VinFast production assumptions (base vs. conservative case)"
    ),
    14: (
        "Contact information:\n"
        "- Tung Ho: tah@allotropepartners.com (technical lead)\n"
        "- Hang Tran: httt@allotropepartners.com (engagement lead)\n\n"
        "Allotrope VC | GTB Vietnam\n\n"
        "Thank the audience for their time and express enthusiasm for potential collaboration.\n\n"
        "Leave-behind: This deck (PDF or PPTX), the synthetic demo dashboard URL, and the "
        "BYOL intake template (if requested)."
    )
}

def add_notes():
    """Add speaker notes to the built deck."""
    print(f"Loading deck from {DECK_PATH}")
    prs = Presentation(str(DECK_PATH))
    
    print(f"Adding speaker notes to {len(SPEAKER_NOTES)} slides")
    for slide_idx, notes_text in SPEAKER_NOTES.items():
        if slide_idx >= len(prs.slides):
            print(f"Warning: Slide {slide_idx} out of range (deck has {len(prs.slides)} slides)")
            continue
        
        slide = prs.slides[slide_idx]
        if not slide.has_notes_slide:
            slide.notes_slide
        
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        
        for i, line in enumerate(notes_text.split('\n')):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
        
        print(f"  Added notes to slide {slide_idx}")
    
    print(f"Saving deck to {DECK_PATH}")
    prs.save(str(DECK_PATH))
    print("Done!")

if __name__ == "__main__":
    add_notes()
