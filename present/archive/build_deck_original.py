#!/usr/bin/env python3
"""Build the Bank-Prospect PACTA/TRISK Pitch Deck from the working copy.

This script applies the 15-slide outline from deck_outline.md to the
reference template, replacing text and images while preserving brand assets.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO_ROOT = Path(__file__).parent.parent
DECK_PATH = REPO_ROOT / "present" / "bank_prospect_deck.pptx"
ASSETS_DIR = REPO_ROOT / "present" / "assets"
OUTPUT_PATH = REPO_ROOT / "present" / "bank_prospect_deck.pptx"

SLIDES_TO_DROP = [13, 14, 17, 18, 19, 20, 21, 22, 23, 24]

SLIDE_MAP = {
    0: {  # Slide 1: Cover (ref slide 0)
        "text": {
            2: "PACTA + TRISK: Transition-Risk Analytics for Vietnamese Banks\nAllotrope VC | GTB Vietnam"
        }
    },
    1: {  # Slide 2: Agenda (ref slide 1)
        "text": {
            1: "Agenda",
            2: "The Problem | Our Approach | Case Study | Next Steps"
        }
    },
    4: {  # Slide 3: The Challenge (ref slide 4)
        "text": {
            1: "The Challenge: Decision 263 & Transition Risk for Vietnamese Banks",
            3: "Decision 263 mandates GHG quotas for thermal power, steel, and cement from 2025. Banks must assess borrower transition risk, but lack analytical tools to quantify portfolio alignment or borrower-level financial stress.",
            4: "Vietnamese banks need an integrated analytics stack that measures portfolio alignment with PDP8/NDC targets, quantifies borrower-level NPV and PD impact under transition scenarios, and produces disclosure-ready outputs for IFRS S2.",
            5: "THE CHALLENGE",
            6: "WHAT BANKS NEED"
        }
    },
    2: {  # Slide 4: Our Solution (ref slide 2)
        "text": {
            1: "Our Solution: Three-Layer Compliance & Risk Stack"
        },
        "table": {
            2: [
                ["Layer 1: PCAF", "Layer 2: PACTA", "Layer 3: TRISK", "Layer 4: IFRS S2"],
                ["Measure", "Align", "Stress", "Disclose"],
                ["GHG inventory", "Portfolio alignment", "Borrower stress", "TCFD/ISSB reporting"],
                ["Scope 1, 2, 3", "PDP8/NDC/NZE", "NPV & PD impact", "Board-ready outputs"]
            ]
        }
    },
    3: {  # Slide 5: PACTA (ref slide 3)
        "text": {
            1: "PACTA: Portfolio Alignment Assessment",
            4: "PACTA matches your loanbook to asset-level company data and compares technology mix against Vietnam's PDP8 and NDC targets. Market-share method for power; SDA for cement and steel.",
            5: "Synthetic MCB portfolio: 43 loans, 25 trillion VND. Power sector shows coal capacity 28% above PDP8 target; renewables buildout 15% below 2030 trajectory.",
            6: "Coverage: Power, Automotive, Cement, Steel, Coal Mining",
            7: "Alignment gaps quantify where your portfolio is off-track",
            8: "Company-level results enable borrower-specific engagement",
            9: "Scenario comparison: PDP8 vs IEA NZE vs Vietnam NDC",
            10: "Outputs feed directly into TRISK stress testing",
            11: "Open-source R packages — no licensing fees",
            14: "1,500+ institutions globally",
            24: "Time to first results: 1–3 months"
        },
        "image": {
            16: "05_vn_power_techmix.png"
        }
    },
    5: {  # Slide 6: TRISK (ref slide 5)
        "text": {
            1: "TRISK: Borrower-Level Transition Stress Testing",
            3: "TRISK quantifies how transition scenarios affect individual borrower creditworthiness using DCF models and Merton-based PD estimation.",
            4: "Power sector: Top coal-heavy borrowers show -12% to -18% NPV change under NZE-aligned stress",
            5: "Cement: BF/BOF mills face -8% to -14% NPV impact under carbon price scenarios",
            6: "Steel: EAF routes outperform BF/BOF by 6–10 percentage points",
            7: "Priority ranking combines alignment gap, NPV change, and exposure weight",
            8: "Sensitivity analysis: shock year, discount rate, market passthrough",
            9: "Results update instantly — no recomputation needed",
            10: "Covers all three Decision 263 sectors",
            12: "Scenario Builder lets you drive the five TRISK levers",
            13: "Exportable: CSV download, HTML reports, interactive dashboard",
            19: "Open-source R package on CRAN"
        },
        "image": {
            8: "trisk_power_npv_change.png"
        }
    },
    8: {  # Slide 7: Case Study Power (ref slide 8)
        "text": {
            0: "Case Study: Vietnam Bank Demo — Power Sector Alignment"
        },
        "image": {
            1: "12_vn_alignment_overview.png"
        },
        "table": {
            2: [
                ["Portfolio", "43 loans, ~$1B USD synthetic"],
                ["Power Mix", "Coal 28%, Gas 12%, Hydro 10%, Solar 8%, Wind 5%"],
                ["PDP8 Target", "Coal phase-down, renewables 3x by 2030"],
                ["Gap", "Coal capacity 28% above target; renewables 15% below"],
                ["Disclaimer", "Illustrative / synthetic data"]
            ]
        }
    },
    9: {  # Slide 8: Case Study TRISK (ref slide 9)
        "text": {
            0: "Case Study: TRISK Borrower Stress Rankings",
            4: "NPV Change by Borrower (Power)",
            5: "PD Change & Coal Stranded Risk",
            6: "TRISK ranks borrowers by financial stress severity under transition scenarios. Coal-heavy names (Vinh Tan, Duyen Hai, Mong Duong) show the highest NPV decline and PD increase. Renewable-focused borrowers (Trung Nam, BIM) benefit from transition tailwinds."
        },
        "image": {
            1: "trisk_power_priority_top10.png",
            2: "trisk_power_pd_change.png",
            3: "13_vn_coal_stranded_risk.png"
        }
    },
    10: {  # Slide 9: Key Findings (ref slide 10)
        "text": {
            0: "Case Study: Key Findings & Takeaways",
            12: "Coal Stranded Risk",
            13: "Renewables Gap",
            14: "Borrower Heterogeneity",
            15: "Decision 263 Readiness",
            16: "Interactive Exploration",
            17: "28% of portfolio coal capacity exceeds PDP8 phase-down trajectory. BOT plants face cliff-edge decline post-PPA expiry (~2035).",
            18: "Renewables buildout is 15% below 2030 target. Solar and wind capacity must triple to align with PDP8.",
            19: "NPV impact ranges from -18% (coal-heavy) to +8% (renewable-focused). One-size engagement fails — borrower-specific data is essential.",
            20: "Three-layer stack (PCAF → PACTA → TRISK) addresses all three Decision 263 requirements: inventory, quotas, and reduction plans.",
            21: "Live dashboard with Scenario Builder lets bank teams explore levers in real time. Downloadable HTML reports and CSV exports."
        }
    },
    11: {  # Slide 10: Compliance Value (ref slide 11)
        "text": {
            1: "Compliance Value: Decision 263 / TCFD / ISSB Alignment",
            5: "Cement & Steel SDA: Emission intensity vs. 2030 NDC targets",
            6: "Coal trajectory vs. PDP8 phase-down pathway"
        },
        "image": {
            2: "10_vn_cement_sda.png",
            3: "11_vn_steel_sda.png",
            4: "06_vn_coal_trajectory.png"
        }
    },
    12: {  # Slide 11: Live Dashboard (ref slide 12)
        "text": {
            1: "Live Dashboard: Interactive Portfolio Exploration",
            6: "PACTA Alignment: Sector-level KPIs and interactive charts",
            7: "TRISK Risk: Borrower stress rankings with sensitivity controls",
            8: "Scenario Builder: Drive the five TRISK levers in real time"
        },
        "image": {
            5: "03_vn_coverage_pie.png"
        }
    },
    6: {  # Slide 12: Engagement Offer (ref slide 6)
        "text": {
            0: "Engagement Offer & 2026 Roadmap",
            3: "Proposed Engagement for 2026:\n\nPhase 1: Data intake and loanbook normalization (BYOL)\nPhase 2: PACTA portfolio alignment with real borrower data\nPhase 3: TRISK stress testing and borrower ranking\nPhase 4: IFRS S2 disclosure package and board-ready report\n\nTimeline: 6–9 months from data receipt to final deliverables.",
            2: "Questions?\n\nWe welcome the opportunity to discuss how this analytics stack can support your bank's transition-risk management and Decision 263 compliance."
        }
    },
    13: {  # Slide 13: Next Steps (ref slide 15, becomes index 13 after dropping)
        "text": {
            0: "Proposed Next Steps",
            3: "Immediate actions:\n\n1. Share loanbook template for BYOL intake validation\n2. Identify Decision 263 borrower subset for priority analysis\n3. Schedule technical workshop with risk and ESG teams\n4. Define scenario preferences (PDP8, NDC, IEA NZE)\n5. Agree on reporting language and disclosure timeline",
            2: "We are ready to begin Phase 1 upon receipt of loanbook data.\n\nAll analytics are open-source with no licensing fees.\n\nSynthetic demo available for immediate internal training."
        }
    },
    14: {  # Slide 14: Questions (ref slide 16, becomes index 14 after dropping)
        "text": {
            0: "Questions & Discussion",
            1: "We welcome questions on methodology, data requirements, implementation timeline, or any aspect of the PACTA + TRISK analytics stack."
        }
    },
    7: {  # Slide 15: Contact (ref slide 7)
        "text": {
            4: "Tung Ho, tah@allotropepartners.com\nHang Tran, httt@allotropepartners.com\nAllotrope VC | GTB Vietnam"
        }
    }
}

def set_shape_text(shape, text):
    """Set text in a shape, preserving formatting where possible."""
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        tf.clear()
        for i, line in enumerate(text.split('\n')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            if tf.paragraphs[0].runs:
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = tf.paragraphs[0].runs[0].font.size if tf.paragraphs[0].runs else Pt(12)
                run.font.name = tf.paragraphs[0].runs[0].font.name if tf.paragraphs[0].runs else "Calibri"

def set_table_content(table, data):
    """Populate a table with data."""
    for i, row_data in enumerate(data):
        if i >= len(table.rows):
            break
        row = table.rows[i]
        for j, cell_text in enumerate(row_data):
            if j >= len(row.cells):
                break
            cell = row.cells[j]
            cell.text = cell_text

def replace_image(shape, image_path):
    """Replace the image in a picture shape."""
    if hasattr(shape, "image"):
        try:
            with open(image_path, 'rb') as f:
                image_data = f.read()
            shape.image._blob = image_data
        except Exception as e:
            print(f"Warning: Could not replace image in shape {shape.shape_id}: {e}")

def drop_slides(prs, slide_indices):
    """Remove slides by index (in reverse order to avoid index shifting)."""
    for idx in sorted(slide_indices, reverse=True):
        if idx < len(prs.slides):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

def build_deck():
    """Build the prospect deck from the working copy."""
    print(f"Loading deck from {DECK_PATH}")
    prs = Presentation(str(DECK_PATH))
    
    print(f"Dropping {len(SLIDES_TO_DROP)} slides: {SLIDES_TO_DROP}")
    drop_slides(prs, SLIDES_TO_DROP)
    
    print(f"Applying changes to {len(SLIDE_MAP)} slides")
    for slide_idx, changes in SLIDE_MAP.items():
        if slide_idx >= len(prs.slides):
            print(f"Warning: Slide {slide_idx} out of range (deck has {len(prs.slides)} slides)")
            continue
        
        slide = prs.slides[slide_idx]
        print(f"  Processing slide {slide_idx}")
        
        if "text" in changes:
            for shape_idx, text in changes["text"].items():
                if shape_idx < len(slide.shapes):
                    shape = slide.shapes[shape_idx]
                    set_shape_text(shape, text)
        
        if "table" in changes:
            for shape_idx, data in changes["table"].items():
                if shape_idx < len(slide.shapes):
                    shape = slide.shapes[shape_idx]
                    if shape.has_table:
                        set_table_content(shape.table, data)
        
        if "image" in changes:
            for shape_idx, image_name in changes["image"].items():
                if shape_idx < len(slide.shapes):
                    shape = slide.shapes[shape_idx]
                    image_path = ASSETS_DIR / image_name
                    if image_path.exists():
                        replace_image(shape, image_path)
                    else:
                        print(f"    Warning: Image not found: {image_path}")
    
    print(f"Saving deck to {OUTPUT_PATH}")
    prs.save(str(OUTPUT_PATH))
    print("Done!")

if __name__ == "__main__":
    build_deck()
