#!/usr/bin/env python3
"""Build the Bank-Prospect PACTA/TRISK Pitch Deck — Claude (present-skill) version.

A clean, comparison build against the opencode/qwen v1+v2 decks. It edits a copy
of the reference template in place by shape index, per the `present` skill, and
fixes the defects observed in the opencode renders:

  * Tables: every cell is written explicitly (full matrix incl. blanks) so no
    stale BESS/GTB rows leak through (slides 2 & 8 in opencode v2).
  * Images: placed with aspect ratio preserved and centred in the slot, so the
    landscape charts are never squashed into tall/narrow boxes (slides 9/11/12).
  * Stale content pictures are explicitly removed before new charts are added.
  * Text panels: copy is sized to fit the original boxes (slide 6 overflow).

Output: present/bank_prospect_deck_claude.pptx  (reference is never modified.)
"""

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

REPO = Path(__file__).parent.parent
REF = REPO / "present" / "ref" / "Allotrope GTB updates_Tara Vietnam_December 2025.pptx"
ASSETS = REPO / "present" / "assets"
OUT = REPO / "present" / "bank_prospect_deck_claude.pptx"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Brand palette (sampled from the reference teal masters).
TEAL = (15, 102, 100)
DARK = (40, 40, 40)
GREY = (90, 90, 90)
RED = (178, 52, 44)
WHITE = (255, 255, 255)

# Reference slides to drop (internal-status / wrong-domain BESS slides).
DROP = [13, 14, 17, 18, 19, 20, 21, 22, 23, 24]


def set_text(shape, text, size=None, bold=None, color=None, font="Calibri",
             align=None, anchor=None):
    """Replace a shape's text with one run per line; suppress inherited bullets."""
    tf = shape.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    # collapse to a single paragraph then rebuild
    for p in list(tf.paragraphs):
        p.clear()
    while len(tf.paragraphs) > 1:
        el = tf.paragraphs[-1]._p
        el.getparent().remove(el)

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.clear()
        run = p.add_run()
        run.text = line
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = RGBColor(*color)
        if font is not None:
            run.font.name = font
        if align is not None:
            p.alignment = align
        # suppress any inherited bullet
        pPr = p._p.find(f"{{{A}}}pPr")
        if pPr is None:
            pPr = etree.SubElement(p._p, f"{{{A}}}pPr")
            p._p.insert(0, pPr)
        if pPr.find(f"{{{A}}}buNone") is None:
            etree.SubElement(pPr, f"{{{A}}}buNone")

    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))


def fill_table(table_shape, data, header=True):
    """Write EVERY cell of the table from `data` (a full row x col matrix).

    Cells beyond `data` are blanked, so stale template rows never leak. Header
    row is rendered bold white-on-teal; body cells are small dark text.
    """
    table = table_shape.table
    nrows, ncols = len(table.rows), len(table.columns)
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            val = ""
            if r < len(data) and c < len(data[r]):
                val = data[r][c]
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Calibri"
            if header and r == 0:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*WHITE)
            else:
                run.font.size = Pt(10)
                run.font.bold = (c == 0)
                run.font.color.rgb = RGBColor(*DARK)


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def fit_image(slide, path, bx, by, bw, bh, valign="center", halign="center"):
    """Add an image fitted inside the (bx,by,bw,bh) box, preserving aspect."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = bw, bw / ar
    if h > bh:
        h, w = bh, bh * ar
    x = bx + (bw - w) / 2 if halign == "center" else (bx if halign == "left" else bx + bw - w)
    y = by + (bh - h) / 2 if valign == "center" else (by if valign == "top" else by + bh - h)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def clear_text_shapes(slide, skip=()):
    """Blank every text frame on the slide (keeps shapes/pictures)."""
    for i, sh in enumerate(slide.shapes):
        if i in skip:
            continue
        if sh.has_table:
            continue
        if sh.has_text_frame:
            sh.text_frame.clear()


def drop_slides(prs, indices):
    for idx in sorted(indices, reverse=True):
        if idx < len(prs.slides):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def build():
    print(f"Loading reference: {REF.name}")
    prs = Presentation(str(REF))
    drop_slides(prs, DROP)
    s = list(prs.slides)
    print(f"  {len(s)} slides after dropping {len(DROP)}")

    # ----- Slide 0 (ref 0): Cover -----
    clear_text_shapes(s[0])
    set_text(s[0].shapes[2],
             "PACTA + TRISK\nTransition-Risk Analytics for Vietnamese Banks",
             size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s[0], "Cover. Pitch a continuation of the GTB Vietnam engagement: an "
                "integrated PACTA + TRISK analytics stack for transition-risk "
                "management and Decision 263 compliance. All figures shown are "
                "illustrative / synthetic.")

    # ----- Slide 1 (ref 1): Agenda -----
    clear_text_shapes(s[1])
    set_text(s[1].shapes[1], "Agenda", size=30, bold=True, color=DARK)
    set_text(s[1].shapes[2],
             "1  The challenge: Decision 263 & transition risk\n"
             "2  Our solution: a three-layer compliance stack\n"
             "3  PACTA portfolio alignment\n"
             "4  TRISK borrower stress testing\n"
             "5  Vietnam-bank case study & live dashboard\n"
             "6  Engagement offer, roadmap & next steps",
             size=13, color=DARK)
    # widen the agenda list box so the lines don't wrap; keep it in the white zone
    s[1].shapes[2].left = Inches(0.6)
    s[1].shapes[2].top = Inches(1.7)
    s[1].shapes[2].width = Inches(7.5)
    s[1].shapes[2].height = Inches(2.6)
    notes(s[1], "Walk the arc: problem -> approach -> methods -> case study -> offer.")

    # ----- Slide 2 (ref 2): Our Solution (3-col, 6-row table w/ col0 merges) -----
    clear_text_shapes(s[2], skip=())
    set_text(s[2].shapes[1], "Our Solution: A Three-Layer Compliance & Risk Stack",
             size=18, bold=True, color=WHITE)
    fill_table(s[2].shapes[2], [
        ["Layer", "Capability", "Output"],
        ["Measure\n& Align", "PCAF financed emissions", "Scope 1/2/3 inventory"],
        ["", "PACTA vs PDP8 / NDC / NZE", "Technology-share alignment gaps"],
        ["Stress\n& Disclose", "TRISK borrower stress", "NPV & PD impact per borrower"],
        ["", "IFRS S2 / TCFD reporting", "Board-ready disclosure pack"],
        ["", "Open-source R stack", "No licensing fees"],
    ])
    notes(s[2], "Four capabilities, one pipeline: measure, align, stress, disclose. "
                "Each layer feeds the next; all open-source.")

    # ----- Slide 3 (ref 3): PACTA -----
    clear_text_shapes(s[3])
    set_text(s[3].shapes[1], "PACTA: Portfolio Alignment Assessment",
             size=20, bold=True, color=WHITE)
    set_text(s[3].shapes[4],
             "Match the loanbook to asset-level company data and compare technology "
             "mix against Vietnam's targets.", size=11, color=DARK)
    set_text(s[3].shapes[5],
             "Synthetic MCB portfolio: 43 loans, 25 trillion VND across power, "
             "automotive, cement, steel and coal mining.", size=10, color=DARK)
    set_text(s[3].shapes[6], "Market-share method for power & automotive", size=10, color=DARK)
    set_text(s[3].shapes[7], "SDA method for cement & steel intensity", size=10, color=DARK)
    set_text(s[3].shapes[8], "Custom PDP8 scenario from Vietnam's PDP8", size=10, color=DARK)
    set_text(s[3].shapes[9], "Company-level results enable borrower-specific engagement",
             size=10, color=DARK)
    set_text(s[3].shapes[10], "Scenario comparison: PDP8 vs IEA NZE vs Vietnam NDC",
             size=10, color=DARK)
    set_text(s[3].shapes[11], "Open-source R packages — no licensing fees", size=10, color=DARK)
    set_text(s[3].shapes[14], "Methodology used by 1,500+ institutions globally",
             size=8, color=GREY)
    set_text(s[3].shapes[16], "Coal 28% above PDP8 target", size=10, bold=True, color=RED)
    set_text(s[3].shapes[17], "Renewables 15% below 2030 trajectory", size=10, bold=True, color=RED)
    set_text(s[3].shapes[24], "Time to first results: 1–3 months", size=8, color=GREY)
    fit_image(s[3], ASSETS / "05_vn_power_techmix.png", 6.55, 1.30, 3.20, 2.25, valign="top")
    notes(s[3], "PACTA answers 'is the portfolio on a Paris-aligned path?' The power "
                "tech-mix chart shows coal above and renewables below the PDP8/NDC path. "
                "Illustrative / synthetic data.")

    # ----- Slide 4 (ref 4): The Challenge -----
    clear_text_shapes(s[4])
    set_text(s[4].shapes[1], "The Challenge: Decision 263 & Transition Risk",
             size=20, bold=True, color=WHITE)
    set_text(s[4].shapes[3],
             "Decision 263/QD-TTg mandates GHG quotas for thermal power, steel and "
             "cement from 2025. Banks must assess borrower transition risk — but lack "
             "tools to quantify portfolio alignment or borrower-level financial stress "
             "under Paris-aligned scenarios.", size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    set_text(s[4].shapes[4],
             "An integrated analytics stack that measures portfolio alignment with "
             "PDP8/NDC targets, quantifies borrower NPV and PD impact under transition "
             "scenarios, and produces disclosure-ready outputs for IFRS S2 and TCFD.",
             size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    set_text(s[4].shapes[5], "THE CHALLENGE", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    set_text(s[4].shapes[6], "WHAT BANKS NEED", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # remove the stale "3 / SLL Support" decorative group last (shifts indices)
    delete_shape(s[4].shapes[2])
    notes(s[4], "Why now: Decision 263 turns transition risk into a compliance "
                "obligation with quotas. Banks need quantification, not just narrative.")

    # ----- Slide 5 (ref 5): TRISK -----
    clear_text_shapes(s[5])
    set_text(s[5].shapes[1], "TRISK: Borrower-Level Transition Stress Testing",
             size=18, bold=True, color=WHITE)
    set_text(s[5].shapes[3],
             "Quantify each borrower's NPV and PD impact under transition scenarios "
             "using DCF and Merton-based credit models.", size=11, color=DARK)
    set_text(s[5].shapes[4], "Power: coal-heavy borrowers −12% to −18% NPV", size=10, color=DARK)
    set_text(s[5].shapes[5], "Cement: BF/BOF mills −8% to −14% NPV", size=10, color=DARK)
    set_text(s[5].shapes[6], "Steel: EAF routes outperform BF/BOF by 6–10pp", size=10, color=DARK)
    set_text(s[5].shapes[7], "Priority ranking: alignment gap + NPV + exposure", size=10, color=DARK)
    set_text(s[5].shapes[8], "Sensitivity: shock year, discount rate, passthrough", size=10, color=DARK)
    set_text(s[5].shapes[9], "Scenario Builder drives five TRISK levers live", size=10, color=DARK)
    set_text(s[5].shapes[10], "Covers all three Decision 263 sectors", size=10, color=DARK)
    set_text(s[5].shapes[12], "Exportable: CSV, HTML reports, dashboard", size=10, color=DARK)
    set_text(s[5].shapes[13], "Results update instantly — no recomputation", size=10, color=DARK)
    set_text(s[5].shapes[19], "Open-source R package (trisk.model)", size=8, color=GREY)
    fit_image(s[5], ASSETS / "trisk_power_npv_change.png", 6.50, 1.55, 3.25, 2.55, valign="top")
    notes(s[5], "TRISK converts alignment gaps into credit impact: NPV change and PD "
                "shift per borrower. Coal names fall hardest; renewables benefit. "
                "Illustrative / synthetic data.")

    # ----- Slide 6 (ref 6): Engagement Offer (fit copy inside the rounded panel) -----
    clear_text_shapes(s[6])
    set_text(s[6].shapes[0], "Engagement Offer & 2026 Roadmap", size=20, bold=True, color=WHITE)
    set_text(s[6].shapes[3],
             "Phase 1  Data intake & loanbook normalisation (BYOL)\n"
             "Phase 2  PACTA portfolio alignment on real borrower data\n"
             "Phase 3  TRISK stress testing & borrower ranking\n"
             "Phase 4  IFRS S2 disclosure pack & board-ready report",
             size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    set_text(s[6].shapes[2],
             "Timeline: 6–9 months from data receipt to final deliverables. "
             "All analytics are open-source with no licensing fees. We welcome the "
             "opportunity to discuss how this stack supports your transition-risk "
             "management and Decision 263 compliance.", size=11, color=DARK,
             anchor=MSO_ANCHOR.TOP)
    notes(s[6], "The ask: a phased 6–9 month engagement starting from a loanbook "
                "intake. Open-source means no licence lock-in.")

    # ----- Slide 7 (ref 7): Contact -----
    clear_text_shapes(s[7])
    set_text(s[7].shapes[4],
             "Tung Ho — tah@allotropepartners.com\n"
             "Hang Tran — httt@allotropepartners.com\n"
             "Allotrope VC  |  Greening the Banks Vietnam",
             size=15, color=TEAL, bold=True)
    notes(s[7], "Contact slide. Offer to scope a workshop with the risk and ESG teams.")

    # ----- Slide 8 (ref 8): Case Study — Power Sector (11-row table, fully filled) -----
    clear_text_shapes(s[8])
    set_text(s[8].shapes[0], "Case Study: Vietnam-Bank Demo — Power Sector Alignment",
             size=16, bold=True, color=WHITE)
    fill_table(s[8].shapes[2], [
        ["Metric", "Value"],
        ["Portfolio", "43 loans · 25 tn VND (synthetic)"],
        ["Sectors", "Power, cement, steel"],
        ["Power mix", "Coal 28% · Gas 12% · Hydro 10%"],
        ["", "Solar 8% · Wind 5%"],
        ["PDP8 target", "Coal phase-down; renewables 3× by 2030"],
        ["Alignment gap", "Coal +28% above; renewables −15% below"],
        ["Method", "PACTA market-share (power)"],
        ["Scenarios", "PDP8 · Vietnam NDC · IEA NZE"],
        ["Output", "Company-level alignment results"],
        ["Data", "Illustrative / synthetic"],
    ])
    old = s[8].shapes[1]
    delete_shape(old)
    fit_image(s[8], ASSETS / "12_vn_alignment_overview.png", 5.55, 1.20, 4.20, 3.40, valign="top")
    notes(s[8], "Worked example on the synthetic MCB loanbook. The alignment-gap chart "
                "ranks technologies by distance from target. Illustrative / synthetic data.")

    # ----- Slide 9 (ref 9): Case Study — TRISK Borrower Rankings (clean 2-chart) -----
    clear_text_shapes(s[9])
    set_text(s[9].shapes[0], "Case Study: TRISK Borrower Stress Rankings",
             size=16, bold=True, color=WHITE)
    # left description column (kept box [6])
    set_text(s[9].shapes[6],
             "TRISK ranks borrowers by financial stress under transition scenarios. "
             "Coal-heavy names (Vinh Tan, Duyen Hai, Mong Duong) show the steepest NPV "
             "decline and PD increase; renewable-focused borrowers (Trung Nam, BIM) "
             "gain transition tailwinds.\n\nIllustrative / synthetic data — not actual "
             "portfolio results.", size=10, color=DARK, anchor=MSO_ANCHOR.TOP)
    # capture caption boxes as objects BEFORE mutating the shape tree
    cap_left, cap_right = s[9].shapes[4], s[9].shapes[5]
    set_text(cap_left, "Priority score — top 10 borrowers", size=9, color=GREY, align=PP_ALIGN.CENTER)
    cap_left.left = Inches(3.25); cap_left.top = Inches(4.55); cap_left.width = Inches(3.20)
    set_text(cap_right, "PD change by borrower", size=9, color=GREY, align=PP_ALIGN.CENTER)
    cap_right.left = Inches(6.55); cap_right.top = Inches(4.55); cap_right.width = Inches(3.20)
    # remove the three original pictures, add two charts side by side on the right
    for idx in (3, 2, 1):
        delete_shape(s[9].shapes[idx])
    fit_image(s[9], ASSETS / "trisk_power_priority_top10.png", 3.25, 1.45, 3.20, 3.05, valign="top")
    fit_image(s[9], ASSETS / "trisk_power_pd_change.png", 6.55, 1.45, 3.20, 3.05, valign="top")
    notes(s[9], "Two views of the same borrowers: priority score and PD change. The "
                "point: one-size engagement fails — you need borrower-level data.")

    # ----- Slide 10 (ref 10): Key Findings (5 boxes, text vertically centred) -----
    clear_text_shapes(s[10])
    set_text(s[10].shapes[0], "Case Study: Key Findings & Takeaways",
             size=16, bold=True, color=WHITE)
    heads = {12: "Coal Stranded Risk", 13: "Renewables Gap", 14: "Borrower Heterogeneity",
             15: "Decision 263 Readiness", 16: "Interactive Exploration"}
    for idx, h in heads.items():
        set_text(s[10].shapes[idx], h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bodies = {
        17: "28% of portfolio coal capacity exceeds the PDP8 phase-down trajectory; "
            "BOT plants face a cliff after PPA expiry (~2035).",
        18: "Renewables buildout is 15% below the 2030 target — solar and wind must "
            "roughly triple to align with PDP8.",
        19: "NPV impact spans −18% (coal-heavy) to +8% (renewable-focused). Borrower-"
            "level data is essential.",
        20: "The PCAF → PACTA → TRISK stack covers all Decision 263 requirements: "
            "inventory, quotas and reduction plans.",
        21: "A live dashboard with Scenario Builder lets bank teams explore levers in "
            "real time, with HTML and CSV exports.",
    }
    for idx, b in bodies.items():
        set_text(s[10].shapes[idx], b, size=9, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    notes(s[10], "Five takeaways. Land the headline: transition risk is concentrated, "
                 "measurable and borrower-specific — and this stack makes it auditable.")

    # ----- Slide 11 (ref 11): Compliance Value (3 charts, aspect-preserved) -----
    clear_text_shapes(s[11])
    set_text(s[11].shapes[1], "Compliance Value: Decision 263 / TCFD / ISSB",
             size=18, bold=True, color=WHITE)
    set_text(s[11].shapes[5],
             "Cement & Steel SDA: emission intensity vs. 2030 NDC targets · synthetic data",
             size=9, color=GREY, align=PP_ALIGN.CENTER)
    set_text(s[11].shapes[6], "Coal trajectory vs. PDP8 phase-down", size=9, color=GREY,
             align=PP_ALIGN.CENTER)
    for idx in (4, 3, 2):
        delete_shape(s[11].shapes[idx])
    # left: coal trajectory; centre: cement SDA; right: steel SDA — all top-aligned
    fit_image(s[11], ASSETS / "06_vn_coal_trajectory.png", 0.10, 1.35, 3.30, 2.10, valign="top")
    fit_image(s[11], ASSETS / "10_vn_cement_sda.png", 3.55, 1.35, 3.35, 2.10, valign="top")
    fit_image(s[11], ASSETS / "11_vn_steel_sda.png", 7.05, 1.35, 2.85, 2.10, valign="top")
    notes(s[11], "Compliance is the strongest 'why now'. The same engine outputs map "
                 "directly to Decision 263, TCFD and IFRS S2 / ISSB. Synthetic data.")

    # ----- Slide 12 (ref 12): Portfolio Analytics montage (4 distinct charts) -----
    clear_text_shapes(s[12])
    set_text(s[12].shapes[1], "Live Dashboard: Portfolio Analytics at a Glance",
             size=18, bold=True, color=WHITE)
    set_text(s[12].shapes[6], "Coverage & alignment", size=9, color=GREY, align=PP_ALIGN.CENTER)
    set_text(s[12].shapes[7], "Borrower NPV impact", size=9, color=GREY, align=PP_ALIGN.CENTER)
    set_text(s[12].shapes[8], "Synthetic data", size=9, color=GREY, align=PP_ALIGN.CENTER)
    for idx in (5, 4, 3, 2):
        delete_shape(s[12].shapes[idx])
    # four distinct outputs (no repeats), each aspect-preserved in its slot
    fit_image(s[12], ASSETS / "03_vn_coverage_pie.png", 0.16, 1.27, 3.29, 3.55, valign="top")
    fit_image(s[12], ASSETS / "12_vn_alignment_overview.png", 3.57, 1.27, 3.39, 2.10, valign="top")
    fit_image(s[12], ASSETS / "trisk_power_npv_change.png", 7.05, 1.25, 2.82, 1.80, valign="top")
    fit_image(s[12], ASSETS / "trisk_power_pd_change.png", 7.05, 3.05, 2.82, 1.80, valign="top")
    notes(s[12], "The dashboard is the leave-behind: bankers explore coverage, "
                 "alignment and borrower stress interactively. Synthetic demo data.")

    # ----- Slide 13 (ref 15): Proposed Next Steps -----
    clear_text_shapes(s[13])
    set_text(s[13].shapes[0], "Proposed Next Steps", size=20, bold=True, color=WHITE)
    set_text(s[13].shapes[3],
             "1.  Share the loanbook template for BYOL intake validation\n"
             "2.  Identify the Decision 263 borrower subset for priority analysis\n"
             "3.  Schedule a technical workshop with risk & ESG teams\n"
             "4.  Agree scenario preferences (PDP8, NDC, IEA NZE)\n"
             "5.  Align on reporting language & disclosure timeline",
             size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    set_text(s[13].shapes[2],
             "Ready to begin Phase 1 on receipt of loanbook data. A synthetic demo is "
             "available now for internal training.", size=11, color=DARK,
             anchor=MSO_ANCHOR.TOP)
    notes(s[13], "Concrete next actions, all low-commitment, to keep momentum after "
                 "the meeting.")

    # ----- Slide 14 (ref 16): Questions -----
    clear_text_shapes(s[14])
    set_text(s[14].shapes[0], "Questions & Discussion", size=26, bold=True, color=TEAL)
    set_text(s[14].shapes[1],
             "We welcome questions on methodology, data requirements, implementation "
             "timeline, or any aspect of the PACTA + TRISK analytics stack.",
             size=13, color=GREY)
    notes(s[14], "Open the floor. Common questions: data needs, accuracy of synthetic "
                 "vs. real, timeline, and how outputs map to Decision 263.")

    print(f"Saving: {OUT.name}")
    prs.save(str(OUT))
    print("Done.")


if __name__ == "__main__":
    build()
